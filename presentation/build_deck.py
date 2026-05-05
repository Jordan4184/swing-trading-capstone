# TODO: fix lightGBM/LightGBM key mismatch - see issue 1

"""
Generate the capstone presentation deck.

Reads from results/*.json and results/*.png to build a 16-slide
presentation matching the Institute of Data template.

Usage:
    python presentation/build_deck.py

Output:
    presentation/capstone_deck.pptx
"""

import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ---------------------------------------------------------------------------
# Design system — Ocean Gradient palette (modified for finance/data feel)
# ---------------------------------------------------------------------------

NAVY = RGBColor(0x14, 0x1B, 0x2D)  # near-black background for title slides
DEEP_BLUE = RGBColor(0x06, 0x5A, 0x82)  # primary
TEAL = RGBColor(0x1C, 0x72, 0x93)  # secondary
ACCENT = RGBColor(0x06, 0xA7, 0x7D)  # green accent (gain/positive)
WARNING = RGBColor(0xE6, 0x39, 0x46)  # red (loss/risk)
LIGHT_BG = RGBColor(0xF7, 0xF9, 0xFB)  # very light gray
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_DARK = RGBColor(0x1A, 0x1A, 0x1A)
TEXT_MUTED = RGBColor(0x55, 0x66, 0x77)

HEADER_FONT = "Calibri"
BODY_FONT = "Calibri"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------


def add_background(slide, color):
    """Set entire slide to a solid background color."""
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    # Send to back
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)
    return bg


def add_text(
    slide,
    text,
    left,
    top,
    width,
    height,
    *,
    font_size=14,
    bold=False,
    color=TEXT_DARK,
    align=PP_ALIGN.LEFT,
    font=BODY_FONT,
    anchor=MSO_ANCHOR.TOP,
):
    """Add a text box. Multi-line text (containing \\n) gets split into separate
    paragraphs so alignment is applied to each line correctly."""
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    tf.vertical_anchor = anchor

    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tx


def add_bullets(
    slide,
    bullets,
    left,
    top,
    width,
    height,
    *,
    font_size=16,
    color=TEXT_DARK,
    font=BODY_FONT,
    line_spacing=1.3,
):
    """Add a bulleted list."""
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)

    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(6)
        run = p.add_run()
        run.text = f"•  {bullet}"
        run.font.name = font
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
    return tx


def add_image_fitted(slide, image_path, left, top, max_width, max_height):
    """Add an image, fitting within max bounds while preserving aspect ratio."""
    pic = slide.shapes.add_picture(str(image_path), left, top)
    # Scale to fit
    w_ratio = max_width / pic.width
    h_ratio = max_height / pic.height
    ratio = min(w_ratio, h_ratio)
    pic.width = int(pic.width * ratio)
    pic.height = int(pic.height * ratio)
    # Center within the box
    pic.left = left + (max_width - pic.width) // 2
    pic.top = top + (max_height - pic.height) // 2
    return pic


def add_slide_header(slide, title, subtitle=None):
    """Standard content slide header."""
    add_text(
        slide,
        title,
        Inches(0.6),
        Inches(0.4),
        Inches(12),
        Inches(0.7),
        font_size=32,
        bold=True,
        color=DEEP_BLUE,
        font=HEADER_FONT,
    )
    if subtitle:
        add_text(
            slide,
            subtitle,
            Inches(0.6),
            Inches(1.05),
            Inches(12),
            Inches(0.4),
            font_size=14,
            color=TEXT_MUTED,
            font=BODY_FONT,
        )


def add_slide_number(slide, n, total=16):
    """Discrete slide number bottom right."""
    add_text(
        slide,
        f"{n} / {total}",
        Inches(12.4),
        Inches(7.0),
        Inches(0.8),
        Inches(0.3),
        font_size=10,
        color=TEXT_MUTED,
        align=PP_ALIGN.RIGHT,
    )


def add_stat_callout(
    slide,
    value,
    label,
    left,
    top,
    width,
    height,
    *,
    value_color=DEEP_BLUE,
    value_size=44,
):
    """Big stat number with small label below."""
    add_text(
        slide,
        value,
        left,
        top,
        width,
        Inches(0.9),
        font_size=value_size,
        bold=True,
        color=value_color,
        font=HEADER_FONT,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        label,
        left,
        top + Inches(0.95),
        width,
        Inches(0.4),
        font_size=12,
        color=TEXT_MUTED,
        align=PP_ALIGN.CENTER,
    )


# ---------------------------------------------------------------------------
# Load data
# ---------------------------------------------------------------------------


def load_metrics(results_dir: Path) -> dict:
    with open(results_dir / "backtest_summary.json") as f:
        backtest = json.load(f)
    with open(results_dir / "model_comparison.json") as f:
        models = json.load(f)
    return {"backtest": backtest, "models": models}


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------


def slide_01_title(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_background(slide, NAVY)

    # Accent dot
    dot = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(0.8), Inches(0.8), Inches(0.25), Inches(0.25)
    )
    dot.fill.solid()
    dot.fill.fore_color.rgb = ACCENT
    dot.line.fill.background()

    add_text(
        slide,
        "CAPSTONE PROJECT — DATA SCIENCE & AI",
        Inches(1.2),
        Inches(0.78),
        Inches(10),
        Inches(0.3),
        font_size=12,
        bold=True,
        color=ACCENT,
        font=BODY_FONT,
    )

    add_text(
        slide,
        "ML-Based Cross-Sectional\nSwing Trading Strategy",
        Inches(0.8),
        Inches(2.2),
        Inches(11.5),
        Inches(2.2),
        font_size=54,
        bold=True,
        color=WHITE,
        font=HEADER_FONT,
    )

    add_text(
        slide,
        "Predicting forward returns with walk-forward machine learning,\n"
        "evaluated against honest benchmarks.",
        Inches(0.8),
        Inches(4.6),
        Inches(11),
        Inches(1),
        font_size=20,
        color=RGBColor(0xCA, 0xDC, 0xFC),
        font=BODY_FONT,
    )

    add_text(
        slide,
        "Jordan Donaldson",
        Inches(0.8),
        Inches(6.4),
        Inches(8),
        Inches(0.4),
        font_size=18,
        bold=True,
        color=WHITE,
    )
    add_text(
        slide,
        "Institute of Data  |  May 2026",
        Inches(0.8),
        Inches(6.8),
        Inches(8),
        Inches(0.4),
        font_size=14,
        color=RGBColor(0xCA, 0xDC, 0xFC),
    )


def slide_02_agenda(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WHITE)
    add_slide_header(slide, "Agenda")

    items = [
        ("01", "Project context", "PDT rule removal & retail algo trading"),
        ("02", "Define", "Business problem and data"),
        ("03", "Design", "EDA findings and pipeline architecture"),
        ("04", "Deliver", "Features, models, and backtest results"),
        ("05", "Conclusions", "What worked, honest caveats, next steps"),
    ]

    y = Inches(1.7)
    for num, head, desc in items:
        # Number badge
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(0.8), y, Inches(0.7), Inches(0.7)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = DEEP_BLUE
        circle.line.fill.background()
        add_text(
            slide,
            num,
            Inches(0.8),
            y,
            Inches(0.7),
            Inches(0.7),
            font_size=18,
            bold=True,
            color=WHITE,
            align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )

        add_text(
            slide,
            head,
            Inches(1.7),
            y - Inches(0.05),
            Inches(11),
            Inches(0.5),
            font_size=22,
            bold=True,
            color=TEXT_DARK,
        )
        add_text(
            slide,
            desc,
            Inches(1.7),
            y + Inches(0.45),
            Inches(11),
            Inches(0.4),
            font_size=14,
            color=TEXT_MUTED,
        )
        y += Inches(1.0)

    add_slide_number(slide, 2)


def slide_03_bio(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WHITE)
    add_slide_header(slide, "About me")

    add_text(
        slide,
        "Jordan Donaldson",
        Inches(0.8),
        Inches(1.7),
        Inches(8),
        Inches(0.5),
        font_size=28,
        bold=True,
        color=DEEP_BLUE,
    )
    add_text(
        slide,
        "Aspiring ML / AI Engineer",
        Inches(0.8),
        Inches(2.25),
        Inches(8),
        Inches(0.4),
        font_size=16,
        color=ACCENT,
        bold=True,
    )

    bullets = [
        "Background in Neuroscience — drawn to systems that learn",
        "Transitioning into ML/AI engineering through the IoD program",
        "Particularly interested in real-time decision systems and AI agents",
        "This capstone combines quantitative methods with applied ML — the kind of work I want to do next",
        "Outside of work: streaming, building, learning in public",
    ]
    add_bullets(
        slide, bullets, Inches(0.8), Inches(3.0), Inches(11.5), Inches(4), font_size=16
    )

    add_slide_number(slide, 3)


def slide_04_context(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WHITE)
    add_slide_header(
        slide,
        "Why this project, why now",
        "Regulatory change opens algorithmic trading to a new audience",
    )

    # Big date callout left
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.9), Inches(4.5), Inches(4.7)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = DEEP_BLUE
    box.line.fill.background()
    box.adjustments[0] = 0.05

    add_text(
        slide,
        "JUNE 4",
        Inches(0.8),
        Inches(2.2),
        Inches(4.5),
        Inches(0.6),
        font_size=18,
        bold=True,
        color=ACCENT,
        align=PP_ALIGN.CENTER,
        font=HEADER_FONT,
    )
    add_text(
        slide,
        "2026",
        Inches(0.8),
        Inches(2.8),
        Inches(4.5),
        Inches(1.6),
        font_size=80,
        bold=True,
        color=WHITE,
        align=PP_ALIGN.CENTER,
        font=HEADER_FONT,
    )
    add_text(
        slide,
        "FINRA removes the\nPattern Day Trader rule",
        Inches(0.8),
        Inches(4.4),
        Inches(4.5),
        Inches(1.0),
        font_size=18,
        color=WHITE,
        align=PP_ALIGN.CENTER,
        font=BODY_FONT,
    )
    add_text(
        slide,
        "$25k minimum equity\nrequirement eliminated",
        Inches(0.8),
        Inches(5.5),
        Inches(4.5),
        Inches(1.0),
        font_size=14,
        color=RGBColor(0xCA, 0xDC, 0xFC),
        align=PP_ALIGN.CENTER,
        font=BODY_FONT,
    )

    # Right side narrative
    add_text(
        slide,
        "What changes",
        Inches(5.8),
        Inches(2.0),
        Inches(7),
        Inches(0.5),
        font_size=20,
        bold=True,
        color=DEEP_BLUE,
    )

    bullets = [
        "Small-account algorithmic strategies become viable for the first time at scale",
        "Retail traders gain access to mechanics previously reserved for $25k+ accounts",
        "Demand emerges for accessible, transparent, ML-augmented trading tools",
        "Question: can retail-grade ML produce defensible alpha with cheap data?",
    ]
    add_bullets(
        slide, bullets, Inches(5.8), Inches(2.6), Inches(7), Inches(4), font_size=15
    )

    add_slide_number(slide, 4)


def slide_05_define_business(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WHITE)
    add_slide_header(slide, "Define — the business problem")

    # Two columns: problem (left), success criteria (right)
    add_text(
        slide,
        "Problem statement",
        Inches(0.8),
        Inches(1.7),
        Inches(6),
        Inches(0.5),
        font_size=20,
        bold=True,
        color=DEEP_BLUE,
    )

    add_text(
        slide,
        "Can a machine-learning system using widely-available retail data\n"
        "produce a swing-trading strategy that delivers risk-adjusted alpha\n"
        "over a passive benchmark, after realistic transaction costs?",
        Inches(0.8),
        Inches(2.3),
        Inches(6),
        Inches(2.2),
        font_size=15,
        color=TEXT_DARK,
    )

    add_text(
        slide,
        "Why this matters",
        Inches(0.8),
        Inches(4.7),
        Inches(6),
        Inches(0.5),
        font_size=20,
        bold=True,
        color=DEEP_BLUE,
    )

    bullets_left = [
        "Validates whether ML adds real value vs. buy-and-hold",
        "Honest framing — most retail strategies underperform",
        "Educational artifact, not a product",
    ]
    add_bullets(
        slide,
        bullets_left,
        Inches(0.8),
        Inches(5.3),
        Inches(6),
        Inches(2),
        font_size=14,
    )

    # Right column — success criteria
    add_text(
        slide,
        "Success criteria",
        Inches(7.3),
        Inches(1.7),
        Inches(5.5),
        Inches(0.5),
        font_size=20,
        bold=True,
        color=DEEP_BLUE,
    )

    success = [
        "Out-of-sample evaluation via walk-forward CV",
        "Realistic transaction costs included (10bps round-trip)",
        "Compared to apples-to-apples benchmark, not just SPY",
        "Reproducible pipeline — single CLI command",
        "Honest discussion of failure modes and risks",
    ]
    add_bullets(
        slide, success, Inches(7.3), Inches(2.3), Inches(5.5), Inches(5), font_size=14
    )

    add_slide_number(slide, 5)


def slide_06_define_data(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WHITE)
    add_slide_header(slide, "Define — the data")

    # Stat row
    stats = [
        ("11", "Tickers"),
        ("8 yrs", "Daily OHLCV"),
        ("22,110", "Raw rows"),
        ("0", "Missing values"),
    ]
    x = Inches(0.8)
    for value, label in stats:
        add_stat_callout(
            slide,
            value,
            label,
            x,
            Inches(1.7),
            Inches(2.9),
            Inches(1.4),
            value_color=DEEP_BLUE,
            value_size=42,
        )
        x += Inches(3.05)

    # Universe + sourcing details
    add_text(
        slide,
        "Universe",
        Inches(0.8),
        Inches(3.6),
        Inches(6),
        Inches(0.4),
        font_size=18,
        bold=True,
        color=DEEP_BLUE,
    )
    add_text(
        slide,
        "AAPL, AMZN, JNJ, JPM, MCD, META, NVDA, PFE, TSLA, UNH  +  SPY (benchmark)",
        Inches(0.8),
        Inches(4.05),
        Inches(11.7),
        Inches(0.5),
        font_size=14,
        color=TEXT_DARK,
    )
    add_text(
        slide,
        "Diverse sectors: tech, consumer, healthcare, financials. "
        "All names traded continuously through the period — no survivorship bias.",
        Inches(0.8),
        Inches(4.5),
        Inches(11.7),
        Inches(0.5),
        font_size=12,
        color=TEXT_MUTED,
    )

    add_text(
        slide,
        "Source & period",
        Inches(0.8),
        Inches(5.2),
        Inches(6),
        Inches(0.4),
        font_size=18,
        bold=True,
        color=DEEP_BLUE,
    )
    add_text(
        slide,
        "Yahoo Finance via yfinance.  2018-01-01 to 2025-12-30.\n"
        "Cached locally as parquet for fast iteration.",
        Inches(0.8),
        Inches(5.65),
        Inches(11.7),
        Inches(1),
        font_size=14,
        color=TEXT_DARK,
    )

    add_slide_number(slide, 6)


def slide_07_design_eda(prs, data, results_dir):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WHITE)
    add_slide_header(
        slide, "Design — EDA findings", "Three findings shaped every downstream choice"
    )

    # Three insight cards
    cards = [
        (
            "Lag-1 mean reversion",
            "Negative autocorrelation across most tickers (SPY -0.14, MCD -0.12). "
            "Prior-day return as a feature.",
        ),
        (
            "Tech cluster correlations",
            "AAPL/SPY pair = 0.77 highest. TSLA/JNJ pair = 0.08 lowest. "
            "Diversification opportunities exist within the universe.",
        ),
        (
            "Cross-sectional > absolute",
            "Forward returns are noisy in absolute terms. "
            "Ranking within the universe is more learnable than predicting magnitude.",
        ),
    ]

    x = Inches(0.8)
    for head, body in cards:
        # Card background
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.8), Inches(4), Inches(2.3)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_BG
        card.line.color.rgb = RGBColor(0xE0, 0xE5, 0xEA)
        card.line.width = Pt(0.75)
        card.adjustments[0] = 0.04

        add_text(
            slide,
            head,
            x + Inches(0.3),
            Inches(2.0),
            Inches(3.6),
            Inches(0.6),
            font_size=16,
            bold=True,
            color=DEEP_BLUE,
        )
        add_text(
            slide,
            body,
            x + Inches(0.3),
            Inches(2.6),
            Inches(3.6),
            Inches(1.5),
            font_size=12,
            color=TEXT_DARK,
        )
        x += Inches(4.15)

    # Correlation matrix image
    corr_path = results_dir / "04_correlation_matrix.png"
    if corr_path.exists():
        add_image_fitted(
            slide, corr_path, Inches(3.3), Inches(4.3), Inches(6.7), Inches(2.9)
        )
        add_text(
            slide,
            "Return correlation matrix (universe)",
            Inches(3.3),
            Inches(7.1),
            Inches(6.7),
            Inches(0.3),
            font_size=10,
            color=TEXT_MUTED,
            align=PP_ALIGN.CENTER,
        )

    add_slide_number(slide, 7)


def slide_08_design_process(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WHITE)
    add_slide_header(
        slide,
        "Design — pipeline architecture",
        "Single CLI command runs the full reproducible pipeline",
    )

    # Process flow: 5 boxes connected by arrows
    steps = [
        ("Data", "Yahoo Finance\nOHLCV + cache", DEEP_BLUE),
        ("Features", "4 families\n(returns, vol,\ntechnicals, market)", TEAL),
        ("Train", "Walk-forward CV\n3 models compared", TEAL),
        ("Backtest", "Non-overlapping\ntop-N w/ costs", TEAL),
        ("Evaluate", "vs. SPY +\nequal-weight", ACCENT),
    ]

    box_w = Inches(2.2)
    box_h = Inches(1.6)
    gap = Inches(0.15)
    total_w = box_w * 5 + gap * 4
    start_x = (SLIDE_W - total_w) // 2
    y = Inches(2.6)

    x = start_x
    for i, (title, body, color) in enumerate(steps):
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, box_w, box_h)
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()
        box.adjustments[0] = 0.08

        add_text(
            slide,
            title,
            x,
            y + Inches(0.15),
            box_w,
            Inches(0.4),
            font_size=16,
            bold=True,
            color=WHITE,
            align=PP_ALIGN.CENTER,
            font=HEADER_FONT,
        )
        add_text(
            slide,
            body,
            x,
            y + Inches(0.6),
            box_w,
            Inches(1.0),
            font_size=11,
            color=WHITE,
            align=PP_ALIGN.CENTER,
        )

        # Arrow to next box (skip last)
        if i < len(steps) - 1:
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                x + box_w + Inches(0.01),
                y + Inches(0.65),
                gap - Inches(0.02),
                Inches(0.3),
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = TEXT_MUTED
            arrow.line.fill.background()

        x += box_w + gap

    # Caption underneath
    add_text(
        slide,
        "python -m src.pipeline --all      runs the entire flow end-to-end. "
        "Each step is independently testable.",
        Inches(0.8),
        Inches(4.6),
        Inches(11.7),
        Inches(0.4),
        font_size=13,
        color=TEXT_MUTED,
        align=PP_ALIGN.CENTER,
        font=BODY_FONT,
    )

    # Below: design choices
    add_text(
        slide,
        "Key design choices",
        Inches(0.8),
        Inches(5.2),
        Inches(11),
        Inches(0.4),
        font_size=18,
        bold=True,
        color=DEEP_BLUE,
    )

    choices = [
        "Walk-forward validation (not random k-fold) — avoids future-data leakage on time series",
        "Cross-sectional ranking target — more learnable than absolute return prediction",
        "Three model classes compared — linear, tree-based ensemble, gradient-boosted",
    ]
    add_bullets(
        slide,
        choices,
        Inches(0.8),
        Inches(5.7),
        Inches(11.7),
        Inches(1.6),
        font_size=13,
    )

    add_slide_number(slide, 8)


def slide_09_deliver_features(prs, data, results_dir):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WHITE)
    add_slide_header(
        slide,
        "Deliver — feature engineering",
        "Four families, each motivated by EDA or financial theory",
    )

    # Four feature family cards in a 2x2 grid
    families = [
        (
            "01",
            "Multi-horizon returns",
            "1d, 5d, 20d, 60d",
            "Captures both short-term mean reversion and longer-horizon momentum",
        ),
        (
            "02",
            "Volatility regime",
            "20d, 60d rolling std",
            "TSLA/NVDA distributions visibly differ from JNJ/PFE",
        ),
        (
            "03",
            "Technical indicators",
            "RSI-14, Bollinger %B",
            "Standard quant signals — cheap to compute, well-studied",
        ),
        (
            "04",
            "Market-relative",
            "Volume ratio, excess return",
            "Isolates idiosyncratic moves from systematic exposure",
        ),
    ]

    positions = [
        (Inches(0.8), Inches(1.8)),
        (Inches(7.0), Inches(1.8)),
        (Inches(0.8), Inches(4.3)),
        (Inches(7.0), Inches(4.3)),
    ]

    for (left, top), (num, title, sub, desc) in zip(positions, families):
        # Number circle
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, left, top, Inches(0.7), Inches(0.7)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = ACCENT
        circle.line.fill.background()
        add_text(
            slide,
            num,
            left,
            top,
            Inches(0.7),
            Inches(0.7),
            font_size=14,
            bold=True,
            color=WHITE,
            align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )

        add_text(
            slide,
            title,
            left + Inches(0.85),
            top - Inches(0.05),
            Inches(5.2),
            Inches(0.4),
            font_size=18,
            bold=True,
            color=DEEP_BLUE,
        )
        add_text(
            slide,
            sub,
            left + Inches(0.85),
            top + Inches(0.4),
            Inches(5.2),
            Inches(0.35),
            font_size=13,
            color=ACCENT,
            bold=True,
            font=BODY_FONT,
        )
        add_text(
            slide,
            desc,
            left + Inches(0.85),
            top + Inches(0.85),
            Inches(5.2),
            Inches(1.2),
            font_size=12,
            color=TEXT_DARK,
        )

    # Bottom note
    add_text(
        slide,
        "Target: binary — 1 if forward-5d return ranks in top 20% of universe on that date.",
        Inches(0.8),
        Inches(6.7),
        Inches(11.7),
        Inches(0.4),
        font_size=13,
        color=TEXT_MUTED,
        align=PP_ALIGN.CENTER,
        bold=True,
    )

    add_slide_number(slide, 9)


def slide_10_deliver_models(prs, data, results_dir):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WHITE)
    add_slide_header(
        slide,
        "Deliver — model comparison",
        "Three classifiers, walk-forward CV, AUC as primary metric",
    )

    models_metrics = data["models"]

    # Three columns for three models
    model_order = ["LogisticRegression", "RandomForest", "lightGBM"]
    display_names = {
        "LogisticRegression": "Logistic Regression",
        "RandomForest": "Random Forest",
        "LightGBM": "LightGBM",
    }
    descriptions = {
        "LogisticRegression": "Linear baseline\nwith StandardScaler",
        "RandomForest": "Non-linear ensemble\nselected as best",
        "LightGBM": "Gradient-boosted\ntrees",
    }

    # Find best by AUC (with tiebreak: RandomForest > LightGBM > LogisticRegression
    # to match models.py behavior — when AUCs tie, prefer the more sophisticated model)
    tiebreak = {"RandomForest": 3, "LightGBM": 2, "LogisticRegression": 1}
    best = max(
        model_order,
        key=lambda m: (
            models_metrics.get(m, {}).get("mean_auc", 0),
            tiebreak.get(m, 0),
        ),
    )

    col_w = Inches(3.8)
    col_gap = Inches(0.2)
    total_w = col_w * 3 + col_gap * 2
    start_x = (SLIDE_W - total_w) // 2
    y = Inches(1.9)

    x = start_x
    for model_key in model_order:
        m = models_metrics.get(model_key, {})
        is_best = model_key == best

        # Card
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, col_w, Inches(4.3)
        )
        card.fill.solid()
        if is_best:
            card.fill.fore_color.rgb = DEEP_BLUE
            card.line.fill.background()
            text_color = WHITE
            label_color = RGBColor(0xCA, 0xDC, 0xFC)
            value_color = ACCENT
        else:
            card.fill.fore_color.rgb = LIGHT_BG
            card.line.color.rgb = RGBColor(0xE0, 0xE5, 0xEA)
            card.line.width = Pt(0.75)
            text_color = TEXT_DARK
            label_color = TEXT_MUTED
            value_color = DEEP_BLUE
        card.adjustments[0] = 0.05

        # "BEST" tag
        if is_best:
            add_text(
                slide,
                "★ BEST BY AUC",
                x + Inches(0.3),
                y + Inches(0.2),
                col_w - Inches(0.6),
                Inches(0.3),
                font_size=11,
                bold=True,
                color=ACCENT,
                align=PP_ALIGN.CENTER,
            )

        # Model name
        add_text(
            slide,
            display_names[model_key],
            x + Inches(0.3),
            y + Inches(0.55),
            col_w - Inches(0.6),
            Inches(0.5),
            font_size=20,
            bold=True,
            color=text_color,
            align=PP_ALIGN.CENTER,
            font=HEADER_FONT,
        )

        add_text(
            slide,
            descriptions[model_key],
            x + Inches(0.3),
            y + Inches(1.1),
            col_w - Inches(0.6),
            Inches(0.7),
            font_size=12,
            color=label_color,
            align=PP_ALIGN.CENTER,
        )

        # Metrics
        metrics_to_show = [
            ("AUC", m.get("mean_auc", 0)),
            ("Accuracy", m.get("mean_accuracy", 0)),
            ("F1", m.get("mean_f1", 0)),
        ]
        my = y + Inches(2.0)
        for lbl, val in metrics_to_show:
            add_text(
                slide,
                f"{val:.3f}",
                x + Inches(0.3),
                my,
                col_w - Inches(0.6),
                Inches(0.5),
                font_size=24,
                bold=True,
                color=value_color,
                align=PP_ALIGN.CENTER,
                font=HEADER_FONT,
            )
            add_text(
                slide,
                lbl,
                x + Inches(0.3),
                my + Inches(0.5),
                col_w - Inches(0.6),
                Inches(0.3),
                font_size=11,
                color=label_color,
                align=PP_ALIGN.CENTER,
            )
            my += Inches(0.75)

        x += col_w + col_gap

    # Insight footnote
    add_text(
        slide,
        "Linear and non-linear models perform similarly — the predictive signal lives in feature engineering, not in complex interactions.",
        Inches(0.8),
        Inches(6.5),
        Inches(11.7),
        Inches(0.5),
        font_size=13,
        color=TEXT_MUTED,
        align=PP_ALIGN.CENTER,
        bold=True,
    )

    add_slide_number(slide, 10)


def slide_11_deliver_equity_curve(prs, data, results_dir):
    """The wow slide — equity curve."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WHITE)
    add_slide_header(
        slide,
        "Deliver — backtest equity curves",
        "6.4 years out-of-sample, $10k initial, 10bps round-trip costs",
    )

    # Big chart taking most of the slide
    chart_path = results_dir / "06_equity_curve.png"
    if chart_path.exists():
        add_image_fitted(
            slide, chart_path, Inches(0.6), Inches(1.6), Inches(8.5), Inches(5.5)
        )

    # Right-side stat callouts
    bt = data["backtest"]
    strategy = bt["strategy"]
    eq = bt["equal_weight"]

    # Strategy
    box1 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.5), Inches(1.8), Inches(3.4), Inches(1.6)
    )
    box1.fill.solid()
    box1.fill.fore_color.rgb = DEEP_BLUE
    box1.line.fill.background()
    box1.adjustments[0] = 0.08

    add_text(
        slide,
        "ML STRATEGY",
        Inches(9.5),
        Inches(1.95),
        Inches(3.4),
        Inches(0.3),
        font_size=11,
        bold=True,
        color=ACCENT,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        f"${10000 * (1 + strategy['total_return']):,.0f}",
        Inches(9.5),
        Inches(2.25),
        Inches(3.4),
        Inches(0.7),
        font_size=32,
        bold=True,
        color=WHITE,
        align=PP_ALIGN.CENTER,
        font=HEADER_FONT,
    )
    add_text(
        slide,
        f"{strategy['annualized_return']*100:.1f}% annualized",
        Inches(9.5),
        Inches(2.95),
        Inches(3.4),
        Inches(0.3),
        font_size=12,
        color=RGBColor(0xCA, 0xDC, 0xFC),
        align=PP_ALIGN.CENTER,
    )

    # Equal-Weight
    box2 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.5), Inches(3.6), Inches(3.4), Inches(1.4)
    )
    box2.fill.solid()
    box2.fill.fore_color.rgb = LIGHT_BG
    box2.line.color.rgb = RGBColor(0xE0, 0xE5, 0xEA)
    box2.adjustments[0] = 0.08

    add_text(
        slide,
        "EQUAL-WEIGHT",
        Inches(9.5),
        Inches(3.75),
        Inches(3.4),
        Inches(0.3),
        font_size=11,
        bold=True,
        color=TEAL,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        f"${10000 * (1 + eq['total_return']):,.0f}",
        Inches(9.5),
        Inches(4.05),
        Inches(3.4),
        Inches(0.6),
        font_size=26,
        bold=True,
        color=DEEP_BLUE,
        align=PP_ALIGN.CENTER,
        font=HEADER_FONT,
    )
    add_text(
        slide,
        f"{eq['annualized_return']*100:.1f}% annualized",
        Inches(9.5),
        Inches(4.65),
        Inches(3.4),
        Inches(0.3),
        font_size=11,
        color=TEXT_MUTED,
        align=PP_ALIGN.CENTER,
    )

    # Alpha box
    box3 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.5), Inches(5.2), Inches(3.4), Inches(1.7)
    )
    box3.fill.solid()
    box3.fill.fore_color.rgb = ACCENT
    box3.line.fill.background()
    box3.adjustments[0] = 0.08

    alpha_pp = (strategy["annualized_return"] - eq["annualized_return"]) * 100
    sharpe_advantage = strategy["sharpe_ratio"] - eq["sharpe_ratio"]

    add_text(
        slide,
        "HONEST ALPHA",
        Inches(9.5),
        Inches(5.35),
        Inches(3.4),
        Inches(0.3),
        font_size=11,
        bold=True,
        color=WHITE,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        f"+{alpha_pp:.2f}pp",
        Inches(9.5),
        Inches(5.65),
        Inches(3.4),
        Inches(0.7),
        font_size=32,
        bold=True,
        color=WHITE,
        align=PP_ALIGN.CENTER,
        font=HEADER_FONT,
    )
    add_text(
        slide,
        f"+{sharpe_advantage:.2f} Sharpe vs. equal-weight",
        Inches(9.5),
        Inches(6.4),
        Inches(3.4),
        Inches(0.4),
        font_size=11,
        color=WHITE,
        align=PP_ALIGN.CENTER,
    )

    add_slide_number(slide, 11)


def slide_12_deliver_performance(prs, data, results_dir):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WHITE)
    add_slide_header(
        slide,
        "Deliver — performance breakdown",
        "Strategy vs. SPY vs. equal-weight universe",
    )

    bt = data["backtest"]

    # Build a metrics table
    metrics = [
        ("Total Return", "total_return", "pct"),
        ("Annualized Return", "annualized_return", "pct"),
        ("Sharpe Ratio", "sharpe_ratio", "num"),
        ("Max Drawdown", "max_drawdown", "pct"),
        ("Hit Rate", "hit_rate", "pct"),
    ]

    table_left = Inches(0.8)
    table_top = Inches(1.9)
    col1_w = Inches(3.5)
    col_w = Inches(2.8)
    row_h = Inches(0.65)

    # Header row
    headers = ["Metric", "ML Strategy", "Equal-Weight", "SPY"]
    header_colors = [DEEP_BLUE, DEEP_BLUE, TEAL, TEXT_MUTED]
    x = table_left
    for i, (h, c) in enumerate(zip(headers, header_colors)):
        w = col1_w if i == 0 else col_w
        # Header background
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, table_top, w, row_h)
        bg.fill.solid()
        bg.fill.fore_color.rgb = c
        bg.line.fill.background()
        add_text(
            slide,
            h,
            x,
            table_top,
            w,
            row_h,
            font_size=15,
            bold=True,
            color=WHITE,
            align=PP_ALIGN.CENTER if i > 0 else PP_ALIGN.LEFT,
            anchor=MSO_ANCHOR.MIDDLE,
            font=HEADER_FONT,
        )
        if i == 0:
            # Add slight left padding via offset
            pass
        x += w

    # Data rows
    keys = [("strategy", DEEP_BLUE), ("equal_weight", TEAL), ("spy", TEXT_MUTED)]
    y = table_top + row_h

    for row_idx, (label, key, fmt) in enumerate(metrics):
        # Alternating row backgrounds
        if row_idx % 2 == 0:
            row_bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, table_left, y, col1_w + col_w * 3, row_h
            )
            row_bg.fill.solid()
            row_bg.fill.fore_color.rgb = LIGHT_BG
            row_bg.line.fill.background()

        # Metric name
        add_text(
            slide,
            label,
            table_left + Inches(0.2),
            y,
            col1_w,
            row_h,
            font_size=14,
            color=TEXT_DARK,
            anchor=MSO_ANCHOR.MIDDLE,
            bold=True,
        )

        # Values
        x = table_left + col1_w
        strategy_val = bt["strategy"].get(key, 0)
        for k, color in keys:
            val = bt[k].get(key, 0)
            if fmt == "pct":
                txt = (
                    f"{val*100:+.2f}%"
                    if key == "max_drawdown" or "return" in key.lower()
                    else f"{val*100:.2f}%"
                )
                if key == "hit_rate":
                    txt = f"{val*100:.1f}%"
            else:
                txt = f"{val:.2f}"

            # Highlight if this is the best in the row (for return/sharpe)
            best_in_row = False
            if key in ("total_return", "annualized_return", "sharpe_ratio", "hit_rate"):
                vals = [bt[kk].get(key, 0) for kk, _ in keys]
                if val == max(vals):
                    best_in_row = True
            elif key == "max_drawdown":
                vals = [bt[kk].get(key, 0) for kk, _ in keys]
                if val == max(vals):  # least negative
                    best_in_row = True

            text_color = ACCENT if best_in_row else color
            add_text(
                slide,
                txt,
                x,
                y,
                col_w,
                row_h,
                font_size=15,
                bold=best_in_row,
                color=text_color,
                align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE,
                font=HEADER_FONT,
            )
            x += col_w

        y += row_h

    # Footnote
    add_text(
        slide,
        "Green = best in row.  6.4-year out-of-sample window, 10bps round-trip costs, $10k initial capital.",
        Inches(0.8),
        Inches(6.0),
        Inches(11.7),
        Inches(0.4),
        font_size=12,
        color=TEXT_MUTED,
        align=PP_ALIGN.CENTER,
        bold=True,
    )

    add_text(
        slide,
        "Strategy delivers higher absolute return AND higher risk-adjusted return than both benchmarks.\n"
        "The cost: deeper max drawdown, reflecting concentration in only 2 names at any time.",
        Inches(0.8),
        Inches(6.5),
        Inches(11.7),
        Inches(0.8),
        font_size=13,
        color=TEXT_DARK,
        align=PP_ALIGN.CENTER,
    )

    add_slide_number(slide, 12)


def slide_13_deliver_risk(prs, data, results_dir):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WHITE)
    add_slide_header(
        slide,
        "Deliver — risk and drawdown",
        "Concentration is the explicit cost of conviction",
    )

    bt = data["backtest"]

    # Drawdown chart on the left
    dd_path = results_dir / "07_drawdown.png"
    if dd_path.exists():
        add_image_fitted(
            slide, dd_path, Inches(0.6), Inches(1.7), Inches(7.5), Inches(4.0)
        )

    # Risk discussion on the right
    add_text(
        slide,
        "What the chart shows",
        Inches(8.4),
        Inches(1.8),
        Inches(4.5),
        Inches(0.4),
        font_size=18,
        bold=True,
        color=DEEP_BLUE,
    )

    risk_bullets = [
        f"Max drawdown: {bt['strategy']['max_drawdown']*100:.1f}%",
        "Largest losses cluster in 2023 — concentrated bets in falling tech",
        "Recovery to new highs by 2024 — strategy is resilient but volatile",
    ]
    add_bullets(
        slide,
        risk_bullets,
        Inches(8.4),
        Inches(2.3),
        Inches(4.5),
        Inches(2),
        font_size=13,
    )

    add_text(
        slide,
        "Honest framing",
        Inches(8.4),
        Inches(4.4),
        Inches(4.5),
        Inches(0.4),
        font_size=18,
        bold=True,
        color=WARNING,
    )

    framing = [
        "Top-2 of 11 = high concentration",
        "Drawdowns are a feature, not a bug",
        "Real money sizing requires risk parity or volatility scaling",
    ]
    add_bullets(
        slide, framing, Inches(8.4), Inches(4.9), Inches(4.5), Inches(2), font_size=13
    )

    # Bottom summary stat strip
    stats = [
        (f"{bt['strategy']['max_drawdown']*100:.1f}%", "Strategy DD", WARNING),
        (f"{bt['equal_weight']['max_drawdown']*100:.1f}%", "Equal-Weight DD", TEAL),
        (f"{bt['spy']['max_drawdown']*100:.1f}%", "SPY DD", TEXT_MUTED),
    ]
    x = Inches(0.8)
    for value, label, color in stats:
        add_stat_callout(
            slide,
            value,
            label,
            x,
            Inches(6.1),
            Inches(2.4),
            Inches(1.0),
            value_color=color,
            value_size=28,
        )
        x += Inches(2.6)

    add_slide_number(slide, 13)


def slide_14_solution(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WHITE)
    add_slide_header(
        slide,
        "End-to-end solution",
        "What I shipped — reproducible, documented, deployable",
    )

    # Two columns: code structure | what each delivers
    add_text(
        slide,
        "Pipeline modules",
        Inches(0.8),
        Inches(1.7),
        Inches(6),
        Inches(0.4),
        font_size=20,
        bold=True,
        color=DEEP_BLUE,
    )

    modules = [
        ("data_loader.py", "Cached OHLCV ingestion"),
        ("features.py", "4 feature families, ~15 signals"),
        ("models.py", "Walk-forward training, 3 algorithms"),
        ("backtest.py", "Non-overlapping trades + 2 benchmarks"),
        ("pipeline.py", "Single-CLI orchestration"),
    ]

    y = Inches(2.3)
    for name, desc in modules:
        # Code-style name in monospace-ish
        add_text(
            slide,
            name,
            Inches(0.8),
            y,
            Inches(2.8),
            Inches(0.4),
            font_size=14,
            bold=True,
            color=TEAL,
            font="Consolas",
        )
        add_text(
            slide,
            desc,
            Inches(3.7),
            y,
            Inches(3.5),
            Inches(0.4),
            font_size=14,
            color=TEXT_DARK,
        )
        y += Inches(0.55)

    # Right: deliverables
    add_text(
        slide,
        "Deliverables",
        Inches(7.8),
        Inches(1.7),
        Inches(5.2),
        Inches(0.4),
        font_size=20,
        bold=True,
        color=DEEP_BLUE,
    )

    deliverables = [
        "Reproducible pipeline (single CLI command)",
        "EDA notebook with 4 published plots",
        "3 model comparison + best-by-AUC selection",
        "Backtest with 2 honest benchmarks + costs",
        "7 result visualizations + JSON metrics",
        "Roadmap doc for next-phase enhancements",
        "Comprehensive README on public GitHub",
    ]
    add_bullets(
        slide,
        deliverables,
        Inches(7.8),
        Inches(2.3),
        Inches(5.2),
        Inches(4.5),
        font_size=13,
    )

    add_text(
        slide,
        "Public repo: github.com/Jordan4184/swing-trading-capstone",
        Inches(0.8),
        Inches(6.7),
        Inches(12),
        Inches(0.4),
        font_size=13,
        color=TEAL,
        align=PP_ALIGN.CENTER,
        bold=True,
    )

    add_slide_number(slide, 14)


def slide_15_summary(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WHITE)
    add_slide_header(slide, "Summary, conclusions, next steps")

    bt = data["backtest"]
    alpha_pp = (
        bt["strategy"]["annualized_return"] - bt["equal_weight"]["annualized_return"]
    ) * 100

    # Three columns: What worked | Honest caveats | Next steps
    headers = [
        ("What worked", ACCENT),
        ("Honest caveats", WARNING),
        ("Next steps", DEEP_BLUE),
    ]

    contents = [
        [
            f"+{alpha_pp:.2f}pp annualized alpha vs. honest benchmark",
            "Walk-forward validation produced reliable estimates",
            "Volatility regime ranks as #1 feature — aligned with theory",
            "Cross-sectional ranking learned despite low signal-to-noise",
        ],
        [
            "Yahoo Finance is an unofficial data source",
            "11-ticker universe is small, may not generalize",
            "2019-2025 was an exceptional bull market for tech",
            "Concentration drawdown of -54% is the cost of conviction",
        ],
        [
            "Crowd Consensus Meter — sentiment as filter layer",
            "Real-time deployment via Polygon + Alpaca paper trading",
            "AI agent for news/filing summarization",
            "Universe expansion from 11 to ~500 names",
        ],
    ]

    col_w = Inches(4.0)
    col_gap = Inches(0.15)
    total_w = col_w * 3 + col_gap * 2
    start_x = (SLIDE_W - total_w) // 2

    x = start_x
    for (head, color), bullets in zip(headers, contents):
        # Top stripe
        stripe = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x, Inches(1.7), col_w, Inches(0.07)
        )
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = color
        stripe.line.fill.background()

        add_text(
            slide,
            head,
            x,
            Inches(1.9),
            col_w,
            Inches(0.5),
            font_size=20,
            bold=True,
            color=color,
            font=HEADER_FONT,
        )

        add_bullets(
            slide,
            bullets,
            x,
            Inches(2.5),
            col_w,
            Inches(4.5),
            font_size=13,
            line_spacing=1.4,
        )

        x += col_w + col_gap

    add_slide_number(slide, 15)


def slide_16_thanks(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, NAVY)

    # Accent dot
    dot = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(0.8), Inches(0.8), Inches(0.25), Inches(0.25)
    )
    dot.fill.solid()
    dot.fill.fore_color.rgb = ACCENT
    dot.line.fill.background()

    add_text(
        slide,
        "QUESTIONS",
        Inches(1.2),
        Inches(0.78),
        Inches(10),
        Inches(0.3),
        font_size=12,
        bold=True,
        color=ACCENT,
        font=BODY_FONT,
    )

    add_text(
        slide,
        "Thank you",
        Inches(0.8),
        Inches(2.5),
        Inches(11.5),
        Inches(1.5),
        font_size=80,
        bold=True,
        color=WHITE,
        font=HEADER_FONT,
    )

    add_text(
        slide,
        "Happy to dig into any part of the pipeline.",
        Inches(0.8),
        Inches(4.2),
        Inches(11.5),
        Inches(0.6),
        font_size=22,
        color=RGBColor(0xCA, 0xDC, 0xFC),
        font=BODY_FONT,
    )

    add_text(
        slide,
        "Jordan Donaldson",
        Inches(0.8),
        Inches(6.3),
        Inches(8),
        Inches(0.4),
        font_size=18,
        bold=True,
        color=WHITE,
    )
    add_text(
        slide,
        "github.com/Jordan4184/swing-trading-capstone   |   itsjordandonaldson@gmail.com",
        Inches(0.8),
        Inches(6.75),
        Inches(11),
        Inches(0.4),
        font_size=13,
        color=RGBColor(0xCA, 0xDC, 0xFC),
    )


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------


def build_deck(project_root: Path, output_path: Path) -> None:
    results_dir = project_root / "results"
    data = load_metrics(results_dir)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    print("Building 16-slide capstone deck...")
    slide_01_title(prs, data)
    slide_02_agenda(prs, data)
    slide_03_bio(prs, data)
    slide_04_context(prs, data)
    slide_05_define_business(prs, data)
    slide_06_define_data(prs, data)
    slide_07_design_eda(prs, data, results_dir)
    slide_08_design_process(prs, data)
    slide_09_deliver_features(prs, data, results_dir)
    slide_10_deliver_models(prs, data, results_dir)
    slide_11_deliver_equity_curve(prs, data, results_dir)
    slide_12_deliver_performance(prs, data, results_dir)
    slide_13_deliver_risk(prs, data, results_dir)
    slide_14_solution(prs, data)
    slide_15_summary(prs, data)
    slide_16_thanks(prs, data)

    output_path.parent.mkdir(exist_ok=True)
    prs.save(str(output_path))
    print(f"Deck saved to {output_path}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    project_root = Path(__file__).parent.parent
    output_path = project_root / "presentation" / "capstone_deck.pptx"
    build_deck(project_root, output_path)
